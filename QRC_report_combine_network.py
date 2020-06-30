# -*- coding: utf-8 -*-
"""
Created on Tue Dec 10 10:57:41 2019

@author: akshay.kale
"""

import pandas as pd
import matplotlib.pyplot as plt
import numpy as np
import warnings
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.dml.color import RGBColor
from PIL import Image, ImageOps
from textwrap import wrap
#from init import QRC_PATH, PPT_PATH, IMG_PATH, DATE_1, DATE_2, NETWORK_PARAMETER_PREPAID, NETWORK_PARAMETER_POSTPAID
from constants import variable

warnings.filterwarnings("ignore")

def create_df(data, network_params, group_no):
    '''Create df as per the requirement'''
    table = pd.pivot_table(data, values ='CNT', index =['SR_SUB_AREA1'], 
                         columns =['CREATED_DATE'], aggfunc = np.sum)
    sorted_table = table.sort_values(table.columns.tolist(), axis=0, ascending=False)
    df = pd.DataFrame(sorted_table.to_records())
    print(df.columns)
    print("-----------------------------------------------------------------------")
    index_no = 1 + (group_no + 1)
    df['before'] = df.iloc[:,1:index_no].sum(axis=1) # Check date range and adjust the index accordingly.
    df['after'] = df.iloc[:, index_no:(index_no+(group_no+1))].sum(axis=1) # Check date range and adjust the index accordingly.
    df = df[['SR_SUB_AREA1', 'before', 'after']]
    idx = df.SR_SUB_AREA1[df.SR_SUB_AREA1.isin(network_params)].index.tolist()
    network_series_sum = df.loc[idx, :].sum() 
    network_df = pd.DataFrame([{'SR_SUB_AREA1': 'Network Related', 'before':int(network_series_sum[1]), 'after': int(network_series_sum[2])}])
    df.drop(df.loc[df['SR_SUB_AREA1'].isin(network_params)].index, inplace=True)
    df = df.sort_values(['before', 'after'], axis=0, ascending=False).head(5)
    df = pd.concat([network_df, df], ignore_index=True, sort =True)
    df = df[['SR_SUB_AREA1', 'before', 'after']]
    df.columns = ['Call Category', variable.DATE_1, variable.DATE_2] #Give name to columns as required 
    df[variable.DATE_1] = df[variable.DATE_1].astype(int) #Change column names.
    df[variable.DATE_2] = df[variable.DATE_2].astype(int) #Change column names.
    print(df)
    return df

def add_text_box(slide, text_data, position, font_size):
    '''Create text box in the slides for headings'''
    left = Inches(position[0])
    top = Inches(position[1])
    width = Inches(position[2])
    height = Inches(position[3])
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = text_data
    p.font.size = Pt(font_size)
    tf.alignment = PP_ALIGN.CENTER


def df_to_table(slide, df, img_name, text_1, text_2, img_title, range_plot, table_heading, img_size, left=None, top=None, width=None,
                height=None, colnames=None, col_formatters=None, rounding=None, name=None):
    
    rows, cols = df.shape
    '''Add a table to a slide'''
    shp = slide.shapes.add_table(rows+2, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    
    '''Merge first row and add heading to that row of table'''
    table = shp.table
    cell = table.cell(0, 0)
    other_cell = table.cell(0,2)
    cell.merge(other_cell)
    cell.text = table_heading
    cell.text_frame.paragraphs[0].font.size = Pt(10)
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    if colnames is None:
        colnames = list(df.columns)

    '''Insert the column names'''
    for col_index, col_name in enumerate(colnames):
        shp.table.cell(1, col_index).text = col_name
        fill = shp.table.cell(1, col_index).fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(31,119,180)
        font = shp.table.cell(1, col_index).text_frame.paragraphs[0].font
        font.size = Pt(10)
        font.color.rgb = RGBColor(255, 255, 255)
        shp.table.cell(1, col_index).text_frame.paragraphs[0].bold = True
        
    m = df.values
    '''Adding data to rows'''
    for row in range(rows):
        for col in range(cols):
            val = m[row, col]
            text = str(val)
            shp.table.cell(row+2, col).text = text
            font = shp.table.cell(row+2, col).text_frame.paragraphs[0].font
            font.size = Pt(10)
        
    if name is not None:
        shp.name = name
    '''Creating plots for count of QRC data.'''    
    plt.rcParams.update({'font.size': 5.9, 'legend.fontsize': 5}) 
    min_count = df[variable.DATE_1].min() #Change the column name as mentioned above
    max_count = df[variable.DATE_1].max() #Change the column name as mentioned above
    df.plot(x='Call Category', y= [variable.DATE_1, variable.DATE_2], kind='bar', width = 0.4, figsize=(6.9,2.35), title=img_title, ylim=(min_count-range_plot[0], max_count+range_plot[1]), rot=0)
    #plt.xticks(wrap = True)
    #plt.tick_params(axis='x', pad=-3)
    
    labels = ['\n'.join(wrap(l,12)) for l in df['Call Category']]
    plt.xticks(range(0, len(df['Call Category'])), labels)
    plt.savefig(variable.IMG_PATH + img_name +'.png', dpi=100)
    #plt.show()

    '''Save image and plot it in ppt'''
    img_path = variable.IMG_PATH + img_name +'.png'
    ImageOps.expand(Image.open(img_path),border=1,fill='black').save(img_path)
    left = Inches(img_size[0])
    top = Inches(img_size[1])
    height = Inches(img_size[2])
    pic = slide.shapes.add_picture(img_path, left, top)    
    
    add_text_box(slide, text_1, [3.5, 0, 0.5, 0.5], 20)
    add_text_box(slide, text_2, [3.5, 2.7, 0.5, 0.5], 16)
    
    '''Adding a straight line in slide'''
    shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0), Inches(0.7), Inches(10), Inches(0.7))
    shape.line.color.rgb = RGBColor(0, 0, 0)
    return shp

def create_slide(df, img_name, text_1, text_2, img_title, range_plot, table_heading):
    '''Creating new slide in existing presentation'''
    pres = Presentation(variable.PPT_PATH)
    blank_slide_layout = pres.slide_layouts[6]
    slide = pres.slides.add_slide(blank_slide_layout)
    table1 = df_to_table(slide, df,img_name, text_1, text_2, img_title, range_plot, table_heading, [0.2, 3.4, 2], 2.5, 0.8, 5, 0.009)
    pres.save(variable.PPT_PATH)
    return table1

# =============================================================================
# data = pd.read_excel(QRC_PATH, sheet_name="Prepaid")
# df = create_df(data, NETWORK_PARAMETER_PREPAID)
# create_slide(df, 'qrc_prepaid', 'QRC Analysis Prepaid', 'Top Call Drivers in Prepaid', 'Prepaid QRC', [12000, 18000], 'QRC Analysis of Prepaid based on CRM data')
# data1 = pd.read_excel(QRC_PATH, sheet_name="Postpaid")
# df1 = create_df(data1, NETWORK_PARAMETER_POSTPAID)
# create_slide(df1, 'qrc_postpaid', 'QRC Analysis Postpaid', 'Top Call Drivers in Postpaid', 'Postpaid QRC', [1000, 1000], 'QRC Analysis of Postpaid based on CRM data')
# 
# print("Created !!!")
# =============================================================================

def qrc_main(start, end, grouping_no, file_path, output_path):
    months = {'01':'Jan', '02':'Feb', '03':'Mar', '04':'Apr', '05':'May', '06':'Jun', '07':'Jul', '08':'Aug', '09':'Sep', '10':'Oct', '11':'Nov', '12':'Dec'}
    date_s = start.split('-')
    date_e = end.split('-')
    month = [months[mo] for mo in months if mo == date_s[1]][0]
    middle_date = int(date_s[2]) + int(grouping_no)
    if grouping_no > 0:
        variable.DATE_1 = date_s[2] + "-" + str(middle_date)  + " " + month
        variable.DATE_2 = "" + str(middle_date + 1) + "-" + date_e[2]  + " " + month
    else:
        variable.DATE_1 = date_s[2] + " " + month
        variable.DATE_2 = date_e[2] + " " + month
    variable.QRC_PATH = file_path
    variable.PPT_PATH = output_path
    
    data = pd.read_excel(variable.QRC_PATH, sheet_name="Prepaid")
    df = create_df(data, variable.NETWORK_PARAMETER_PREPAID, grouping_no)
    create_slide(df, 'qrc_prepaid', 'QRC Analysis Prepaid', 'Top Call Drivers in Prepaid', 'Prepaid QRC', [12000, 130000], 'QRC Analysis of Prepaid based on CRM data')
    data1 = pd.read_excel(variable.QRC_PATH, sheet_name="Postpaid")
    df1 = create_df(data1, variable.NETWORK_PARAMETER_POSTPAID, grouping_no)
    create_slide(df1, 'qrc_postpaid', 'QRC Analysis Postpaid', 'Top Call Drivers in Postpaid', 'Postpaid QRC', [1000, 6700], 'QRC Analysis of Postpaid based on CRM data')
    print("Created QRC!!!")
    return "Success"