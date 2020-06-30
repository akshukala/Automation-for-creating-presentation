# -*- coding: utf-8 -*-
"""
Created on Fri Nov  1 11:20:25 2019

@author: akshay.kale
"""

import pandas as pd
import matplotlib.pyplot as plt
import warnings
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.dml.color import RGBColor
from PIL import Image, ImageOps
#from init import IVR_PATH, IMG_PATH, PPT_PATH, ST_DATE_1, ST_DATE_2
from constants import variable
warnings.filterwarnings("ignore")


def plot_chart(file_name, sheet_name, chart_title, offered_range, helpline):
    '''Export data from excel'''
    data = pd.read_excel(variable.IVR_PATH + file_name +".xlsx", sheet_name=sheet_name, skiprows=33, nrows=23, usecols=range(0,32))
    index_start_no = int(variable.ST_DATE_1.split('-')[2])
    index_end_no = int(variable.ST_DATE_2.split('-')[2]) + 1
    df = data.iloc[:, index_start_no:index_end_no] # Adjust the data index as per requirement.
    print (df)
    offered_df = df.sum(axis=0, skipna=True)
    ''' File path '''
    at_data = pd.read_excel(variable.IVR_PATH + file_name +".xlsx", sheet_name=sheet_name, skiprows=153, nrows=23, usecols=range(0,32))
    at_data_temp = at_data.iloc[:, index_start_no:index_end_no] # Adjust the data index as per requirement.
    at_df = at_data_temp.sum(axis=0, skipna=True)
    '''Creating a dataframe for plotting a graph'''
    final_df = offered_df.to_frame()
    final_df.columns = ['Offered_Calls']
    final_df['AT'] = at_df
    final_df = final_df.astype(int)
    final_df['time']= final_df.index
    final_df['ATPercent'] = (final_df['AT'] / final_df['Offered_Calls'])*100
    final_df.reset_index(inplace = True)
    final_df = final_df.drop(['time'], axis=1)
    max_at_percent = final_df['ATPercent'].max()
    min_at_percent = final_df['ATPercent'].min()
    max_offered_calls = final_df['Offered_Calls'].max() 
    print(final_df)
    
    '''chart preparing code'''
    width = 0.35
    plt.rcParams.update({'font.size': 5.9, 'legend.fontsize': 5})
    plt.tight_layout()
    final_df[['Offered_Calls','AT']].plot(kind='bar', width = width, figsize=(3.5,2.1), title=chart_title, ylim=(0,(max_offered_calls + offered_range)))
    final_df['ATPercent'].plot(secondary_y=True, color='brown', linestyle = 'dashed', linewidth = 2, marker='o', markerfacecolor='brown')
    
    '''Dual Y axis'''
    ax = plt.gca()
    ax.set_ylim([(min_at_percent - 3),(max_at_percent + 3)])
    plt.xlim([-width, len(final_df['Offered_Calls'])-width])
    #ax.xaxis.set_major_locator(plt.MaxNLocator(4))
    ax.set_xticklabels(final_df['index'].dt.strftime('%d-%b-%y'))
    '''Labels for line chart'''
    for i,j in final_df.ATPercent.items():
        ax.annotate(str(round(j,2)) + "%", xy=(i, j+0.3), ha='left')
    ax.legend(loc='upper left', ncol=2)
    plt.savefig(variable.IMG_PATH +sheet_name.replace(' ', '') + helpline +'.png', dpi=100)
    plt.tight_layout()
    #plt.show()
    return "Chart Created"

def create_slide(image_names, flag, heading, position_array):
    if flag == 1:
        '''When creating new presentation'''
        prs = Presentation()
    else:
        '''Adding slides to existing file'''
        p = open(variable.PPT_PATH, 'rb')
        prs = Presentation(p)
    
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    '''Slide Heading'''
    left = Inches(1.8)
    top = Inches(0)
    width = Inches(0.5)
    height = Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = heading
    p.font.size = Pt(23)
    tf.alignment = PP_ALIGN.CENTER
    
    '''Setting charts on slide'''
    for i,name in enumerate(image_names):
        img_path = variable.IMG_PATH +image_names[i]+'.png'
        ImageOps.expand(Image.open(img_path),border=1,fill='black').save(img_path)
        left = Inches(position_array[i][0])
        top = Inches(position_array[i][1])
        height = Inches(position_array[i][2])
        pic = slide.shapes.add_picture(img_path, left, top)
    
    shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0), Inches(0.8), Inches(10), Inches(0.8))
    shape.line.color.rgb = RGBColor(0, 0, 0)
    
    prs.save(variable.PPT_PATH)
    print("Created")


    
# =============================================================================
# overall_198 = plot_chart('198_Performance', 'Trend 198 IVR', '198 Overall Trend', 210000, '198')
# postpaid_198 = plot_chart('198_Performance', 'Trend Postpaid', '198 Postpaid Trend', 6000, '198')
# prepaid_198 = plot_chart('198_Performance', 'Trend Prepaid', '198 Prepaid Trend', 180000, '198')
# prospect_198 = plot_chart('198_Performance', 'Trend Prospect', '198 Prospect Trend', 10000, '198')
# image_names = ['Trend198IVR198', 'TrendPostpaid198', 'TrendPrepaid198', 'TrendProspect198']
# position_array = [(0.05, 0.85, 2), (5.05, 0.85, 2), (0.05, 3.85, 2), (5.05, 3.85, 2)]
# create_slide(image_names, 1, "198 Detailed analysis ( Segregation of Data)", position_array)
# 
# 
# overall_12345 = plot_chart('12345_Performance', 'Trend 12345 IVR', '12345 Overall Trend', 125000, '12345')
# postpaid_12345 = plot_chart('12345_Performance', 'Trend Postpaid', '12345 Postpaid Trend', 9000, '12345')
# prepaid_12345 = plot_chart('12345_Performance', 'Trend Prepaid', '12345 Prepaid Trend', 90000, '12345')
# prospect_12345 = plot_chart('12345_Performance', 'Trend Prospect', '12345 Prospect Trend', 8000, '12345')
# image_names = ['Trend12345IVR12345', 'TrendPostpaid12345', 'TrendPrepaid12345', 'TrendProspect12345']
# position_array = [(0.05, 0.85, 2), (5.05, 0.85, 2), (0.05, 3.85, 2), (5.05, 3.85, 2)]
# create_slide(image_names, 2, "12345 Detailed analysis ( Segregation of Data)", position_array)
# 
# postpaid_multimodal_12345 = plot_chart('12345_Performance', 'Trend Postpaid Multi Modal', 'Trend of Multimodal Postpaid', 1200, '12345')
# prepaid_multimodal_12345 = plot_chart('12345_Performance', 'Trend Prepaid Multi Modal', 'Trend of Multimodal Prepaid', 30000, '12345')
# image_names = ['TrendPostpaidMultiModal12345', 'TrendPrepaidMultiModal12345']
# position_array = [(0.05, 0.85, 2), (5.05, 0.85, 2)]
# create_slide(image_names, 2, "12345 Detailed analysis ( Segregation of Data)", position_array)
# =============================================================================


def main_1(start, end, ivr_path, output_path):
    variable.ST_DATE_1 = start
    variable.ST_DATE_2 = end
    variable.IVR_PATH = ivr_path
    variable.PPT_PATH = output_path
    
    overall_198 = plot_chart('198_Performance', 'Trend 198 IVR', '198 Overall Trend', 210000, '198')
    postpaid_198 = plot_chart('198_Performance', 'Trend Postpaid', '198 Postpaid Trend', 6000, '198')
    prepaid_198 = plot_chart('198_Performance', 'Trend Prepaid', '198 Prepaid Trend', 180000, '198')
    prospect_198 = plot_chart('198_Performance', 'Trend Prospect', '198 Prospect Trend', 10000, '198')
    image_names = ['Trend198IVR198', 'TrendPostpaid198', 'TrendPrepaid198', 'TrendProspect198']
    position_array = [(0.05, 0.85, 2), (5.05, 0.85, 2), (0.05, 3.85, 2), (5.05, 3.85, 2)]
    create_slide(image_names, 1, "198 Detailed analysis ( Segregation of Data)", position_array)
    
    overall_12345 = plot_chart('12345_Performance', 'Trend 12345 IVR', '12345 Overall Trend', 125000, '12345')
    postpaid_12345 = plot_chart('12345_Performance', 'Trend Postpaid', '12345 Postpaid Trend', 9000, '12345')
    prepaid_12345 = plot_chart('12345_Performance', 'Trend Prepaid', '12345 Prepaid Trend', 90000, '12345')
    prospect_12345 = plot_chart('12345_Performance', 'Trend Prospect', '12345 Prospect Trend', 8000, '12345')
    image_names = ['Trend12345IVR12345', 'TrendPostpaid12345', 'TrendPrepaid12345', 'TrendProspect12345']
    position_array = [(0.05, 0.85, 2), (5.05, 0.85, 2), (0.05, 3.85, 2), (5.05, 3.85, 2)]
    create_slide(image_names, 2, "12345 Detailed analysis ( Segregation of Data)", position_array)
    
    postpaid_multimodal_12345 = plot_chart('12345_Performance', 'Trend Postpaid Multi Modal', 'Trend of Multimodal Postpaid', 1200, '12345')
    prepaid_multimodal_12345 = plot_chart('12345_Performance', 'Trend Prepaid Multi Modal', 'Trend of Multimodal Prepaid', 30000, '12345')
    image_names = ['TrendPostpaidMultiModal12345', 'TrendPrepaidMultiModal12345']
    position_array = [(0.05, 0.85, 2), (5.05, 0.85, 2)]
    create_slide(image_names, 2, "12345 Detailed analysis ( Segregation of Data)", position_array)
    
    return "Success from IVR"