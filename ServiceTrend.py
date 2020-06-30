# -*- coding: utf-8 -*-
"""
Created on Tue Nov  5 16:43:48 2019

@author: akshay.kale
"""


import pandas as pd
import matplotlib.pyplot as plt
import warnings
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.dml.color import RGBColor
from PIL import Image, ImageOps
#from init import  PPT_PATH, IMG_PATH, ST_PATH, ST_DATE_1, ST_DATE_2
from constants import variable

warnings.filterwarnings("ignore")

def df_to_table(slide, df, tbl_text, img_name, graph_text, img_size, left=None, top=None, width=None,
                height=None, colnames=None, col_formatters=None, rounding=None, name=None):
    

    rows, cols = df.shape
    '''Add a table to a slide'''
    shp = slide.shapes.add_table(rows+2, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    
    '''Merge first row and add heading to that row of table'''
    table = shp.table
    cell = table.cell(0, 0)
    other_cell = table.cell(0,3)
    cell.merge(other_cell)
    cell.text = tbl_text
    cell.text_frame.paragraphs[0].font.size = Pt(10)
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    if colnames is None:
        colnames = list(df.columns)

    '''Insert the column names'''
    for col_index, col_name in enumerate(colnames):
        shp.table.cell(1,col_index).text = col_name
        fill = shp.table.cell(1, col_index).fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(31,119,180)
        font = shp.table.cell(1, col_index).text_frame.paragraphs[0].font
        font.size = Pt(10)
        font.color.rgb = RGBColor(255, 255, 255)
        shp.table.cell(1, col_index).text_frame.paragraphs[0].bold = True
        
    m = df.values
    '''Adding data to rows'''
    for row in range(rows):
        for col in range(cols):
            val = m[row, col]
            text = str(val)
            shp.table.cell(row+2, col).text = text
            font = shp.table.cell(row+2, col).text_frame.paragraphs[0].font
            font.size = Pt(10)
        
    if name is not None:
        shp.name = name
    '''Creating plots for AT Percent'''    
    plt.rcParams.update({'font.size': 5.9, 'legend.fontsize': 5}) 
    graph_df = df.drop(['Offered', 'Service %'], axis=1)
    min_at_percent = df['Agent Transfer %'].min()
    max_at_percent = df['Agent Transfer %'].max()
    graph_df.plot(x = 'Dates', y = 'Agent Transfer %', kind='bar', width = 0.4, figsize=(3.5,1.65), title=graph_text, ylim=((min_at_percent-1),(max_at_percent+1)), rot=0)
    for i,j in enumerate(graph_df['Agent Transfer %']):
        plt.annotate(str(round(j,2)) + "%", xy=(i, j+0.1), ha='center')
    plt.savefig(variable.IMG_PATH + img_name +'.png', dpi=100)
    #plt.show()

    '''Save image and plot it in ppt'''
    img_path = variable.IMG_PATH + img_name +'.png'
    ImageOps.expand(Image.open(img_path),border=1,fill='black').save(img_path)
    left = Inches(img_size[0])
    top = Inches(img_size[1])
    height = Inches(img_size[2])
    pic = slide.shapes.add_picture(img_path, left, top)    
    
    left_header = Inches(2.5)
    top_header = Inches(0)
    width_header = Inches(0.5)
    height_header = Inches(0.5)
    txBox = slide.shapes.add_textbox(left_header, top_header, width_header, height_header)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = 'Offered Calls Trend Analysis'
    p.font.size = Pt(20)
    tf.alignment = PP_ALIGN.CENTER
    
    shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0), Inches(0.7), Inches(10), Inches(0.7))
    shape.line.color.rgb = RGBColor(0, 0, 0)    
    
    return shp


def create_plot(slide, df, img_title, img_name):
    plt.rcParams.update({'font.size': 7, 'legend.fontsize': 6}) 
    min_at_percent = df['Transfer %'].min()
    max_at_percent = df['Transfer %'].max()
    df.plot(x = 'Dates', y = 'Transfer %', kind='bar', width = 0.4, figsize=(5,2), title=img_title, ylim=((min_at_percent-2),(max_at_percent+5)), rot=0)
    for i,j in enumerate(df['Transfer %']):
        plt.annotate(str(round(j,2)) + "%", xy=(i, j+0.1), ha='center')
    plt.savefig(variable.IMG_PATH + img_name +'.png', dpi=100)
    #plt.show()

    '''Save image and plot it in ppt'''
    img_path = variable.IMG_PATH + img_name +'.png'
    ImageOps.expand(Image.open(img_path),border=1,fill='black').save(img_path)
    left = Inches(1.4)
    top = Inches(4.3)
    height = Inches(2)
    pic = slide.shapes.add_picture(img_path, left, top)
    
    left_header = Inches(2.5)
    top_header = Inches(0)
    width_header = Inches(0.5)
    height_header = Inches(0.5)
    txBox = slide.shapes.add_textbox(left_header, top_header, width_header, height_header)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = 'Observation from Service Trend Report'
    p.font.size = Pt(20)
    tf.alignment = PP_ALIGN.CENTER
    
    shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0), Inches(0.7), Inches(10), Inches(0.7))
    shape.line.color.rgb = RGBColor(0, 0, 0)

# =============================================================================
# '''Export data from excel'''
# 
# data = pd.read_excel(variable.ST_PATH, sheet_name="Trend Sheet", skiprows=3, nrows=32, usecols=range(1,9))
# 
# data = data.drop(['Unnamed: 5'], axis=1)
#     
# 
# '''Get 198 and 12345 data from sheet'''
# data = data.dropna()
# df_12345 = data.iloc[:, 0:4]
# df_198 = data.iloc[:, 4:]
# df_198['Dates'] = df_12345['Dates']
# df_198.columns = ['Offered', 'Service %', 'Agent Transfer %', 'Dates']
# 
# '''Getting df according to date'''
# df1 = df_12345.loc[(df_12345['Dates'] >= variable.ST_DATE_1) & (df_12345['Dates'] <= variable.ST_DATE_2)]
# df1['Dates'] = df1['Dates'].dt.strftime('%d-%b-%y')
# df1['Offered'] = df1['Offered'].astype(int)
# print (df1)
# df2 = df_198.loc[(df_198['Dates'] >= variable.ST_DATE_1) & (df_198['Dates'] <= variable.ST_DATE_2)]
# df2['Dates'] = df2['Dates'].dt.strftime('%d-%b-%y')
# df2['Offered'] = df2['Offered'].astype(int)
# df2 = df2[['Dates', 'Offered', 'Service %', 'Agent Transfer %']]
# print (df2)
# 
# '''Creating new slide in existing presentation'''
# pres = Presentation(variable.PPT_PATH)
# blank_slide_layout = pres.slide_layouts[6]
# slide = pres.slides.add_slide(blank_slide_layout)
# table1 = df_to_table(slide, df1, '12345 IVR New Flow', 'AT_12345', '12345 Agent Transfer %', [0.05,4.9,2],  0.01, 3.1, 5, 0.009)
# table2 = df_to_table(slide, df2, '198 IVR New Flow', 'AT_198', '198 Agent Transfer %', [5.05, 4.9, 2], 5, 3.1, 5, 0.009)
# 
# data1 = pd.read_excel(variable.ST_PATH, sheet_name="Trend Sheet", skiprows=3, nrows=32, usecols=[1, 19, 20])
# df_at = data1.drop([data1.columns[1]], axis=1)
# df_at = df_at.dropna()
# df_at['Transfer %'] = round(df_at['Transfer %'] * 100, 2) 
# df3 = df_at.loc[(df_at['Dates'] >= variable.ST_DATE_1) & (df_at['Dates'] <= variable.ST_DATE_2)]
# df3['Dates'] = df3['Dates'].dt.strftime('%d-%b-%y')
# blank_slide_layout1 = pres.slide_layouts[6]
# slide1 = pres.slides.add_slide(blank_slide_layout1)
# plot = create_plot(slide1, df3, 'Agent Transfer %', 'AT_Summary')
# print(df3)
# pres.save(variable.PPT_PATH)
# =============================================================================


def service_trend_main(start, end, file_path, output_path):
    variable.ST_DATE_1 = start
    variable.ST_DATE_2 = end
    variable.ST_PATH = file_path
    variable.PPT_PATH = output_path
    
    '''Export data from excel'''

    data = pd.read_excel(variable.ST_PATH, sheet_name="Trend Sheet", skiprows=3, nrows=32, usecols=range(1,9))
    
    data = data.drop(['Unnamed: 5'], axis=1)
        
    
    '''Get 198 and 12345 data from sheet'''
    data = data.dropna()
    df_12345 = data.iloc[:, 0:4]
    df_198 = data.iloc[:, 4:]
    df_198['Dates'] = df_12345['Dates']
    df_198.columns = ['Offered', 'Service %', 'Agent Transfer %', 'Dates']
    
    '''Getting df according to date'''
    df1 = df_12345.loc[(df_12345['Dates'] >= variable.ST_DATE_1) & (df_12345['Dates'] <= variable.ST_DATE_2)]
    df1['Dates'] = df1['Dates'].dt.strftime('%d-%b-%y')
    df1['Offered'] = df1['Offered'].astype(int)
    print (df1)
    df2 = df_198.loc[(df_198['Dates'] >= variable.ST_DATE_1) & (df_198['Dates'] <= variable.ST_DATE_2)]
    df2['Dates'] = df2['Dates'].dt.strftime('%d-%b-%y')
    df2['Offered'] = df2['Offered'].astype(int)
    df2 = df2[['Dates', 'Offered', 'Service %', 'Agent Transfer %']]
    print (df2)
    
    '''Creating new slide in existing presentation'''
    p = open(variable.PPT_PATH, 'rb')
    pres = Presentation(p)
    blank_slide_layout = pres.slide_layouts[6]
    slide = pres.slides.add_slide(blank_slide_layout)
    table1 = df_to_table(slide, df1, '12345 IVR New Flow', 'AT_12345', '12345 Agent Transfer %', [0.05,4.9,2],  0.01, 3.1, 5, 0.009)
    table2 = df_to_table(slide, df2, '198 IVR New Flow', 'AT_198', '198 Agent Transfer %', [5.05, 4.9, 2], 5, 3.1, 5, 0.009)
    
    data1 = pd.read_excel(variable.ST_PATH, sheet_name="Trend Sheet", skiprows=3, nrows=32, usecols=[1, 19, 20])
    df_at = data1.drop([data1.columns[1]], axis=1)
    df_at = df_at.dropna()
    df_at['Transfer %'] = round(df_at['Transfer %'] * 100, 2) 
    df3 = df_at.loc[(df_at['Dates'] >= variable.ST_DATE_1) & (df_at['Dates'] <= variable.ST_DATE_2)]
    df3['Dates'] = df3['Dates'].dt.strftime('%d-%b-%y')
    blank_slide_layout1 = pres.slide_layouts[6]
    slide1 = pres.slides.add_slide(blank_slide_layout1)
    plot = create_plot(slide1, df3, 'Agent Transfer %', 'AT_Summary')
    print(df3)
    pres.save(variable.PPT_PATH)
    
    return "Success from Service Trend"