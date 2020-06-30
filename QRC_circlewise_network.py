# -*- coding: utf-8 -*-
"""
Created on Tue Dec 10 13:03:14 2019

@author: akshay.kale
"""

import pandas as pd
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.dml.color import RGBColor
from PIL import Image, ImageOps
#from init import QRC_PATH, PPT_PATH, IMG_PATH, DATE_1, DATE_2, NETWORK_PARAMETER_PREPAID, NETWORK_PARAMETER_POSTPAID
from constants import variable

def create_df(data, network_params, group_no):
    '''Create df as per the requirement'''
    table = pd.pivot_table(data, values ='CNT', index =['SR_SUB_AREA1'], 
                         columns =['CREATED_DATE'], aggfunc = np.sum)
    sorted_table = table.sort_values(table.columns.tolist(), axis=0, ascending=False)
    flattened = pd.DataFrame(sorted_table.to_records()).head(6)
    index_no = 1 + (group_no + 1)
    flattened['before'] = flattened.iloc[:,1:index_no].sum(axis=1) # Check date range and adjust the index accordingly.
    flattened['after'] = flattened.iloc[:, index_no:(index_no+(group_no+1))].sum(axis=1) # Check date range and adjust the index accordingly.
    df = flattened[['SR_SUB_AREA1', 'before', 'after']]
    
    idx = df.SR_SUB_AREA1[df.SR_SUB_AREA1.isin(network_params)].index.tolist()
    network_series_sum = df.loc[idx, :].sum() 
    network_df = pd.DataFrame([{'SR_SUB_AREA1': 'Network Related', 'before':int(network_series_sum[1]), 'after': int(network_series_sum[2])}])
    df.drop(df.loc[df['SR_SUB_AREA1'].isin(network_params)].index, inplace=True)
    df = df.sort_values(['before', 'after'], axis=0, ascending=False).head(5)
    df = pd.concat([network_df, df], ignore_index=True, sort =True)
    df = df[['SR_SUB_AREA1', 'before', 'after']]
    
    df = df.sort_values(['before', 'after'], axis=0, ascending=False)
    df.columns = ['SR_SUB_AREA1', variable.DATE_1, variable.DATE_2] #Give name to columns as required 
    df[variable.DATE_1] = df[variable.DATE_1].astype(int) #Change column names.
    df[variable.DATE_2] = df[variable.DATE_2].astype(int) #Change column names.
    print(df)
    top_area = df['SR_SUB_AREA1'][0]
    print(top_area)
    circlewise_table = data[data.SR_SUB_AREA1.isin(network_params)].pivot_table(values ='CNT', index =['X_DIV_NAME'], 
                         columns =['CREATED_DATE'], aggfunc = np.sum)
    sorted_table_1 = circlewise_table.sort_values(circlewise_table.columns.tolist(), axis=0, ascending=False)
    total_sr_area =  sorted_table_1.sum(axis = 0, skipna = True).sum()
    flattened_1 = pd.DataFrame(sorted_table_1.to_records()).head(6)
    flattened_1['before'] = flattened_1.iloc[:,1:index_no].sum(axis=1) # Check date range and adjust the index accordingly.
    flattened_1['after'] = flattened_1.iloc[:, index_no:(index_no+(group_no+1))].sum(axis=1) # Check date range and adjust the index accordingly.
    df_1 = flattened_1[['X_DIV_NAME', 'before', 'after']]
    df_1 = df_1.sort_values(['after', 'before'], axis=0, ascending=False)
    df_1.columns = ['CIRCLE', variable.DATE_1, variable.DATE_2] #Give name to columns as required
    df_1[variable.DATE_1] = df_1[variable.DATE_1].astype(int) #Change column names.
    df_1[variable.DATE_2] = df_1[variable.DATE_2].astype(int) #Change column names.
    print(df_1)
    return [df, df_1, total_sr_area]

def add_text_box(slide, text_data, position, font_size):
    '''Create text box in the slides for headings'''
    left = Inches(position[0])
    top = Inches(position[1])
    width = Inches(position[2])
    height = Inches(position[3])
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = text_data
    p.font.size = Pt(font_size)
    tf.alignment = PP_ALIGN.CENTER


def df_to_table(slide, df, table_heading, img_title, range_plot, img_name, img_size, text_1, 
                left=None, top=None, width=None, height=None, colnames=None,
                col_formatters=None, rounding=None, name=None):
    total_sr = 0 
    if len(df) > 2:
        total_sr = df.pop(2)
    #for i, frame in enumerate(df):
    rows, cols = df[1].shape
    '''Add a table to a slide'''
    shp = slide.shapes.add_table(rows+2, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    
    '''Merge first row and add heading to that row of table'''
    table = shp.table
    cell = table.cell(0, 0)
    other_cell = table.cell(0,2)
    cell.merge(other_cell)
    cell.text = table_heading
    cell.text_frame.paragraphs[0].font.size = Pt(10)
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    if colnames is None:
        colnames = list(df[1].columns)

    '''Insert the column names'''
    for col_index, col_name in enumerate(colnames):
        shp.table.cell(1,col_index).text = col_name
        fill = shp.table.cell(1, col_index).fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(31,119,180)
        font = shp.table.cell(1, col_index).text_frame.paragraphs[0].font
        font.size = Pt(10)
        font.color.rgb = RGBColor(255, 255, 255)
        shp.table.cell(1, col_index).text_frame.paragraphs[0].bold = True
        
    m = df[1].values
    '''Adding data to rows'''
    for row in range(rows):
        for col in range(cols):
            val = m[row, col]
            text = str(val)
            shp.table.cell(row+2, col).text = text
            font = shp.table.cell(row+2, col).text_frame.paragraphs[0].font
            font.size = Pt(10)
        
    if name is not None:
        shp.name = name
        
    '''Creating plots for count of QRC data.'''    
    plt.rcParams.update({'font.size': 7, 'legend.fontsize': 7})
    col_names = df[1].columns.tolist()[1:]
    min_at_percent = df[1][col_names[0]].min()
    max_at_percent = df[1][col_names[1]].max()
    
    #df[1]['percent'] = (df[1][col_names[1]] / total_sr) * 100
    df[1].plot(x = 'CIRCLE', y = col_names, kind='bar', width = 0.4, figsize=(6.9,2), title=img_title, ylim=((min_at_percent-range_plot[0]),(max_at_percent+range_plot[1])), rot=0)
# =============================================================================
#     for i,j in enumerate(df[1]['percent']):
#         plt.text(i, (j + df[1][col_names[1]][i]+ 500), str(round(j,2)) + "%", fontsize=6)
# =============================================================================
    plt.savefig(variable.IMG_PATH + img_name +'.png', dpi=100)
    #plt.show()

    '''Save image and plot it in ppt'''
    img_path = variable.IMG_PATH + img_name +'.png'
    ImageOps.expand(Image.open(img_path),border=1,fill='black').save(img_path)
    left = Inches(img_size[0])
    top = Inches(img_size[1])
    height = Inches(img_size[2])
    pic = slide.shapes.add_picture(img_path, left, top)

    
    add_text_box(slide, text_1 , [3.5, 0, 0.5, 0.5], 20)
    #add_text_box(slide, text_2, [3.5, 3.2, 0.5, 0.5], 16)
    
    '''Adding a straight line in slide'''
    shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0), Inches(0.7), Inches(10), Inches(0.7))
    shape.line.color.rgb = RGBColor(0, 0, 0)
    return shp


def create_slide(df, table_heading, img_title, range_plot, img_name, img_size, text_1):
    '''Creating new slide in existing presentation'''
    pres = Presentation(variable.PPT_PATH)
    blank_slide_layout = pres.slide_layouts[6]
    slide = pres.slides.add_slide(blank_slide_layout)
    table1 = df_to_table(slide, df, table_heading, img_title, range_plot, img_name, img_size, text_1, 2.6, 0.8, 5, 0.0009)
    pres.save(variable.PPT_PATH)
    return table1

# =============================================================================
# data = pd.read_excel(variable.QRC_PATH, sheet_name="Prepaid")
# df = create_df(data, variable.NETWORK_PARAMETER_PREPAID)
# print (df)
# create_slide(df,  'Circle wise count for ' + df[0]['SR_SUB_AREA1'][0],
#              'Circle wise count for ' + df[0]['SR_SUB_AREA1'][0], [3000, 3500],
#              'circlewise_pre', [0.2,3,2.5], 'Prepaid: ' + df[0]['SR_SUB_AREA1'][0] + ' QRC')
# data1 = pd.read_excel(variable.QRC_PATH, sheet_name="Postpaid")
# df1 = create_df(data1, variable.NETWORK_PARAMETER_POSTPAID)
# create_slide(df1, 'Circle wise count for ' + df1[0]['SR_SUB_AREA1'][0],
#              'Circle wise count for ' + df1[0]['SR_SUB_AREA1'][0], [150, 200],
#              'circlewise_pre', [0.2,3,2.5], 'Postpaid: ' + df1[0]['SR_SUB_AREA1'][0] + ' QRC')
# print("Created")
# =============================================================================

def circlewise_network_main(start, end, grouping_no, output_path):
    months = {'01':'Jan', '02':'Feb', '03':'Mar', '04':'Apr', '05':'May', '06':'Jun', '07':'Jul', '08':'Aug', '09':'Sep', '10':'Oct', '11':'Nov', '12':'Dec'}
    date_s = start.split('-')
    date_e = end.split('-')
    month = [months[mo] for mo in months if mo == date_s[1]][0]
    middle_date = int(date_s[2]) + int(grouping_no)
    if grouping_no > 0:
        variable.DATE_1 = date_s[2] + "-" + str(middle_date)  + " " + month
        variable.DATE_2 = "" + str(middle_date + 1) + "-" + date_e[2]  + " " + month
    else:
        variable.DATE_1 = date_s[2] + " " + month
        variable.DATE_2 = date_e[2] + " " + month
    variable.PPT_PATH = output_path
        
    data = pd.read_excel(variable.QRC_PATH, sheet_name="Prepaid")
    df = create_df(data, variable.NETWORK_PARAMETER_PREPAID, grouping_no)
    print (df)
    create_slide(df,  'Circle wise count for ' + df[0]['SR_SUB_AREA1'][0],
                 'Circle wise count for ' + df[0]['SR_SUB_AREA1'][0], [3000, 3500],
                 'circlewise_pre', [0.2,3,2.5], 'Prepaid: ' + df[0]['SR_SUB_AREA1'][0] + ' QRC')
    data1 = pd.read_excel(variable.QRC_PATH, sheet_name="Postpaid")
    df1 = create_df(data1, variable.NETWORK_PARAMETER_POSTPAID, grouping_no)
    create_slide(df1, 'Circle wise count for ' + df1[0]['SR_SUB_AREA1'][0],
                 'Circle wise count for ' + df1[0]['SR_SUB_AREA1'][0], [250, 1200],
                 'circlewise_pre', [0.2,3,2.5], 'Postpaid: ' + df1[0]['SR_SUB_AREA1'][0] + ' QRC')
    return "Successfully Created circle wise plots."